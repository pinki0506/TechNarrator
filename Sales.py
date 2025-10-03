import pandas as pd
import os
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# ----------------- Setup paths -----------------
BASE_DIR = r"C:\Users\ps186070\TechNarrator"
OUTPUT_DIR = os.path.join(BASE_DIR, "Output")
os.makedirs(OUTPUT_DIR, exist_ok=True)

HEALTH_JSON = os.path.join(BASE_DIR, "health_messages.json")
PPT_FILE = os.path.join(OUTPUT_DIR, "Health_Report.pptx")
CONSOLIDATED_FILE = os.path.join(OUTPUT_DIR, "Consolidated_Final_Health.csv")
COMMENTS_FILE = os.path.join(OUTPUT_DIR, "Consolidated_With_Comments.csv")

# ----------------- Helper functions -----------------
SEVERITY_ORDER = {"Good":1, "Average":2, "Degraded":3, "Critical":4, "N/A":0, "N/A or Invalid Value":0}

def worst_status(status_list):
    return max(status_list, key=lambda x: SEVERITY_ORDER.get(x, 0))

def assess(value, thresholds):
    try:
        v = float(str(value).replace('%','').strip())
    except:
        return "N/A"
    for limit, label in thresholds:
        if v < limit: return label
    return thresholds[-1][1]

def split_text(text, max_len=100):
    """Split long text into lines to fit PPT"""
    words = text.split()
    lines, current = [], ""
    for w in words:
        if len(current) + len(w) + 1 <= max_len:
            current += " " + w if current else w
        else:
            lines.append(current)
            current = w
    if current:
        lines.append(current)
    return lines

# ----------------- Load health messages -----------------
with open(HEALTH_JSON, 'r') as f:
    health_messages = json.load(f)

# ----------------- Process raw CSVs -----------------
def process_awt(file):
    df = pd.read_csv(file)
    df.columns = df.columns.str.strip()
    df['Health_99'] = df['AWT86_99Pct'].apply(lambda x: assess(x, [(10,"Good"),(15,"Average"),(20,"Degraded"),(999,"Critical")]))
    df['Health_100'] = df['AWT_100Pct'].apply(lambda x: assess(x, [(5,"Average"),(10,"Degraded"),(999,"Critical")]))
    df['AWT_Final_Health_Status'] = df.apply(lambda r: worst_status([r['Health_99'], r['Health_100']]), axis=1)
    return df[['ProjectName','AWT_Final_Health_Status']]

def process_flow(file):
    df = pd.read_csv(file)
    df.columns = df.columns.str.strip()
    df['FLOW_Final_Health_Status'] = df.apply(lambda r: worst_status([
        assess(r['FCTime_30_60_secs'], [(10,"Good"),(15,"Average"),(20,"Degraded"),(99999,"Critical")]),
        assess(r['FCTime_1_3_mins'], [(15,"Average"),(20,"Degraded"),(99999,"Critical")]),
        assess(r['FCTime_3_5_mins'], [(20,"Degraded"),(25,"Degraded"),(99999,"Critical")]),
        assess(r['FCTime_5_mins_plus'], [(10,"Degraded"),(99999,"Critical")])
    ]), axis=1)
    return df[['ProjectName','FLOW_Final_Health_Status']]

def process_delay(file):
    df = pd.read_csv(file)
    df.columns = df.columns.str.strip()
    df['Delay_Final_Health_Status'] = df.apply(lambda r: worst_status([
        assess(r['Delay_1_5_min'], [(1000,"Good"),(1500,"Average"),(2500,"Degraded"),(999999,"Critical")]),
        assess(r['Delay_5_10_min'], [(100,"Good"),(200,"Average"),(500,"Degraded"),(999999,"Critical")]),
        assess(r['Delay_10_30_min'], [(50,"Good"),(100,"Average"),(200,"Degraded"),(999999,"Critical")]),
        assess(r['Delay_30_60_min'], [(20,"Good"),(50,"Average"),(100,"Degraded"),(999999,"Critical")]),
        assess(r['Delay_60_min_plus'], [(10,"Average"),(999999,"Critical")])
    ]), axis=1)
    return df[['ProjectName','Delay_Final_Health_Status']]

# ----------------- Merge all -----------------
try:
    awt = process_awt(os.path.join(BASE_DIR,"AWT.csv"))
    flow = process_flow(os.path.join(BASE_DIR,"FlowControl.csv"))
    delay = process_delay(os.path.join(BASE_DIR,"Delay.csv"))
    df_final = awt.merge(flow, on="ProjectName", how="outer").merge(delay, on="ProjectName", how="outer")
    df_final.to_csv(CONSOLIDATED_FILE, index=False)
    print(f"✅ Consolidated health CSV saved at {CONSOLIDATED_FILE}")
except Exception as e:
    print(f"Error processing raw CSVs: {e}")

# ----------------- Add comments -----------------
def generate_comments(row):
    awt_msg = health_messages['AWT'].get(row['AWT_Final_Health_Status'], "AWT status info")
    flow_msg = health_messages['FLOW'].get(row['FLOW_Final_Health_Status'], "FLOW status info")
    delay_msg = health_messages['DELAY'].get(row['Delay_Final_Health_Status'], "DELAY status info")
    return f"AWT: {awt_msg} | FLOW: {flow_msg} | DELAY: {delay_msg}"

df_final['Comments'] = df_final.apply(generate_comments, axis=1)
df_final.to_csv(COMMENTS_FILE, index=False)
print(f"✅ Comments added to {COMMENTS_FILE}")

# ----------------- Create PPT -----------------
colors = {"Good": RGBColor(0,128,0), "Average": RGBColor(255,165,0),
          "Degraded": RGBColor(255,165,0), "Critical": RGBColor(255,0,0), "N/A": RGBColor(128,128,128)}

prs = Presentation()
for _, row in df_final.iterrows():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Project name centered
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = row['ProjectName']
    run.font.bold = True
    run.font.size = Pt(36)
    
    # Comments box
    comments_box = slide.shapes.add_textbox(Inches(1), Inches(1.7), Inches(8), Inches(5))
    tf = comments_box.text_frame
    tf.clear()
    tf.word_wrap = True
    
    # Split comments so each AWT, FLOW, DELAY on new line
    for part in str(row['Comments']).split('|'):
        part = part.strip()
        if ':' in part:
            prefix, text = part.split(':', 1)
            lines = split_text(text.strip(), max_len=100)
            # Each prefix starts on new paragraph
            for i, line in enumerate(lines):
                p = tf.add_paragraph()
                p.alignment = PP_ALIGN.LEFT
                run1 = p.add_run()
                if i == 0:
                    run1.text = prefix + ": "
                    # Set severity color
                    if "AWT" in prefix.upper():
                        sev = row['AWT_Final_Health_Status']
                    elif "FLOW" in prefix.upper():
                        sev = row['FLOW_Final_Health_Status']
                    else:
                        sev = row['Delay_Final_Health_Status']
                    run1.font.color.rgb = colors.get(sev, RGBColor(0,0,0))
                    run1.font.bold = True
                run2 = p.add_run()
                run2.text = line
        else:
            lines = split_text(part.strip(), max_len=100)
            for line in lines:
                p = tf.add_paragraph()
                p.alignment = PP_ALIGN.LEFT
                p.text = line

prs.save(PPT_FILE)
print(f"✅ PPT saved at {PPT_FILE}")
